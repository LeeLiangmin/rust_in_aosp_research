# -*- coding: utf-8 -*-
"""Generate insight deck: Rust in Android/AOSP."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.oxml.ns import qn
import copy

# ---------- palette ----------
INK = RGBColor(0x1F, 0x2A, 0x37)        # near-black slate
BODY = RGBColor(0x37, 0x42, 0x4F)
MUTED = RGBColor(0x6B, 0x75, 0x80)
ACCENT = RGBColor(0xB7, 0x41, 0x0E)     # rust orange
ACCENT2 = RGBColor(0x14, 0x5A, 0x6E)    # deep teal
LIGHT = RGBColor(0xF4, 0xF1, 0xEC)      # warm paper
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set_font(run, size=14, bold=False, color=BODY, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", FONT)


def _fill_para(p, text, size, color, bold_all=False):
    """mini-markdown: **bold** segments."""
    parts = text.split("**")
    for i, seg in enumerate(parts):
        if not seg:
            continue
        r = p.add_run()
        r.text = seg
        _set_font(r, size=size, bold=bold_all or (i % 2 == 1), color=color)


def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def rect(slide, l, t, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def add_footer(slide, note, page):
    tb = textbox(slide, Inches(0.55), Inches(7.02), Inches(11.2), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    _fill_para(p, "数据来源：" + note, 9, MUTED)
    tb2 = textbox(slide, Inches(12.35), Inches(7.02), Inches(0.7), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    _fill_para(p2, str(page), 10, MUTED, bold_all=True)


def content_slide(kicker, title, page, title_size=26):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, Inches(0.16), SH, ACCENT)
    tb = textbox(s, Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.45))
    _fill_para(tb.text_frame.paragraphs[0], kicker, 13, ACCENT, bold_all=True)
    tb = textbox(s, Inches(0.55), Inches(0.72), Inches(12.3), Inches(1.15))
    _fill_para(tb.text_frame.paragraphs[0], title, title_size, INK, bold_all=True)
    rect(s, Inches(0.57), Inches(1.78), Inches(1.35), Pt(3.2), ACCENT)
    return s


def bullets(slide, items, l=Inches(0.6), t=Inches(2.0), w=Inches(12.1), h=Inches(4.8), size=15, gap=10):
    tb = textbox(slide, l, t, w, h)
    tf = tb.text_frame
    first = True
    for text, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap if lvl == 0 else 5)
        p.line_spacing = 1.18
        mark = "▍ " if lvl == 0 else "– "
        r = p.add_run()
        r.text = "    " * lvl + mark
        _set_font(r, size=size if lvl == 0 else size - 1.5, bold=(lvl == 0),
                  color=ACCENT if lvl == 0 else MUTED)
        _fill_para(p, text, size if lvl == 0 else size - 1.5, BODY if lvl else INK)
    return tb


# ============ S1 cover ============
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(6.9), SW, Inches(0.6), ACCENT)
tb = textbox(s, Inches(0.9), Inches(1.1), Inches(11.5), Inches(0.5))
_fill_para(tb.text_frame.paragraphs[0], "洞察报告 · 2026-08", 14, RGBColor(0xE8, 0xB4, 0x9B), bold_all=True)
tb = textbox(s, Inches(0.9), Inches(1.9), Inches(11.6), Inches(2.6))
tf = tb.text_frame
_fill_para(tf.paragraphs[0], "Rust 在 Android 的七年：", 40, WHITE, bold_all=True)
p = tf.add_paragraph()
_fill_para(p, "一场方法论的胜利", 40, RGBColor(0xF0, 0x8A, 0x4B), bold_all=True)
tb = textbox(s, Inches(0.9), Inches(4.35), Inches(11.4), Inches(1.2))
tf = tb.text_frame
_fill_para(tf.paragraphs[0],
           "内存安全漏洞占比 76%（2019）→ <20%（2025），绝对数量 223 → 年化 36。",
           18, RGBColor(0xC9, 0xD2, 0xDA))
p = tf.add_paragraph()
_fill_para(p, "改变结果的不是一门语言，而是一套可复制的渐进改造方法。", 18, RGBColor(0xC9, 0xD2, 0xDA))
tb = textbox(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.5))
_fill_para(tb.text_frame.paragraphs[0],
           "数据底座：Google 官方博客 13 篇逐句回源 · AOSP main 分支 40+ 仓库克隆实测 · 第三方全树统计（AOSP 14/16/17）",
           11, RGBColor(0x9A, 0xA5, 0xB0))

# ============ S2 核心观点 ============
s = content_slide("EXECUTIVE SUMMARY", "四个核心观点", 2)
cards = [
    ("01  胜利的不是 Rust，是方法",
     "“只换增量、不动存量”+ 漏洞随年龄指数衰减（半衰期），让攻击面在存量不动的情况下持续收缩。76%→<20% 是这套数学的必然结果。"),
    ("02  源头预防第一次被证明优于被动缓解",
     "三代策略（修补→缓解→预防）里只有预防让曲线拐头：占比 76%→35%→24%→<20%，绝对数量 223→85→年化 36，远低于 70% 的行业常态。"),
    ("03  收益已从安全外溢到效率，采用开始自我加速",
     "Rust 变更评审快 25%、返工少 20%、回滚率低 4 倍；还省内存省线程。官方原话：“the safer path is now also the faster one.”"),
    ("04  可复制的是路径，不是技术栈",
     "量化互操作可行性（87–90% 可对接）→ 风险驱动选点 → unsafe 收拢（仅 4%）→ 培训体系化 → 用 DORA 指标闭环证明价值。"),
]
top = Inches(2.05)
for i, (h, b) in enumerate(cards):
    card_t = top + Inches(1.22) * i
    rect(s, Inches(0.6), card_t, Inches(12.15), Inches(1.08), LIGHT)
    tb = textbox(s, Inches(0.85), card_t + Inches(0.07), Inches(11.7), Inches(0.42))
    _fill_para(tb.text_frame.paragraphs[0], h, 15, ACCENT, bold_all=True)
    tb = textbox(s, Inches(0.85), card_t + Inches(0.47), Inches(11.7), Inches(0.58))
    _fill_para(tb.text_frame.paragraphs[0], b, 12, BODY)
add_footer(s, "Google Online Security Blog 2019–2026（13 篇原文回源核对）；本报告实测", 2)

# ============ S3 问题定性 ============
s = content_slide("背景 · 为什么必须动手", "内存安全不是普通 bug 类别，而是 Android 漏洞的结构性来源", 3)
bullets(s, [
    ("2019 年官方盘点：UAF、整数溢出、越界读写合计占当年漏洞的 **90%**，其中边界类占 34%；后续追溯口径：内存安全占 **76%**", 0),
    ("这些漏洞“更严重、更易远程触达、更易被实际利用”（2022 原文）——占 36% 的漏洞贡献了大比例的高危案例", 1),
    ("三代应对策略，前两代都触到天花板", 0),
    ("第 1 代 被动修补：打地鼠，修复本身还会引入新漏洞", 1),
    ("第 2 代 运行时缓解（ASLR/sanitizer/CFI/沙箱）：只抬高利用成本；沙箱本身吃内存、加 IPC 延迟，漏洞密度高时可被串联打穿", 1),
    ("第 3 代 源头预防（内存安全语言）：直接消灭漏洞类别——2019 年 Android 团队据此定方向", 1),
    ("**观点：当一个漏洞类别占七成、且缓解手段边际递减，出路只剩一条——换掉生产工具本身**", 0),
], t=Inches(2.0), size=14.5)
add_footer(s, "security.googleblog.com 2019-05《Queue the Hardening Enhancements》、2021-04、2022-12（均已回源）", 3)

# ============ S4 洞察1 方法论 ============
s = content_slide("洞察 ① · 方法论", "漏洞会“衰老”：换掉增量，攻击面就会指数收缩", 4)
bullets(s, [
    ("官方的内部发现（2021）：约 **50%** 的内存安全 bug 诞生于不足一年的新代码", 0),
    ("2024 年形式化为“漏洞半衰期”模型：**5 年陈代码的漏洞密度比新代码低 3.4–7.4 倍**（Usenix Security 2022 研究 + Android/Chromium 观测）", 0),
    ("于是策略变得反直觉但成立：", 0),
    ("**不重写**数千万行存量——越老的代码越安全，重写收益递减还引入新 bug", 1),
    ("只要求**新代码**默认内存安全语言——新代码安全 + 老代码自然衰减 = 总量指数下降", 1),
    ("2024 年实测与模型吻合：多数代码仍“不安全”，但漏洞持续下降（官方原图结论）", 1),
    ("**观点：安全改造的关键决策不是“用什么语言”，而是“改哪部分代码”——这是本案例最可迁移的一课**", 0),
], t=Inches(2.0), size=14.5)
add_footer(s, "security.googleblog.com 2021-04、2024-09《Eliminating Memory Safety Vulnerabilities at the Source》（均已回源）", 4)

# ============ S5 洞察2 数据拐头（图） ============
s = content_slide("洞察 ② · 数据", "七年曲线：源头预防让漏洞占比拐头向下", 5)
chart_data = CategoryChartData()
chart_data.categories = ["2019", "2022", "2024", "2025"]
chart_data.add_series("内存安全漏洞占比 %", (76, 35, 24, 19))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(2.0), Inches(5.9), Inches(4.3), chart_data)
ch = gf.chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = "占比：76% → 35% → 24% → <20%"
for r in ch.chart_title.text_frame.paragraphs[0].runs:
    _set_font(r, 13, True, INK)
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0"%"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(12)
plot.data_labels.font.bold = True
ser = ch.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = ACCENT

chart_data2 = CategoryChartData()
chart_data2.categories = ["2019", "2022", "2024(年化)"]
chart_data2.add_series("内存安全漏洞绝对数量", (223, 85, 36))
gf2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.8), Inches(2.0), Inches(5.9), Inches(4.3), chart_data2)
ch2 = gf2.chart
ch2.has_legend = False
ch2.has_title = True
ch2.chart_title.text_frame.text = "绝对数量：223 → 85 → 36"
for r in ch2.chart_title.text_frame.paragraphs[0].runs:
    _set_font(r, 13, True, INK)
plot2 = ch2.plots[0]
plot2.has_data_labels = True
plot2.data_labels.font.size = Pt(12)
plot2.data_labels.font.bold = True
ser2 = ch2.series[0]
ser2.format.fill.solid()
ser2.format.fill.fore_color.rgb = ACCENT2
bullets(s, [
    ("三个“首次”：2022 内存安全漏洞**首次不再是多数**；2024 降至 24%、**远低于 70% 行业常态**；2025 **首次低于 20%**", 0),
    ("同期安全公告月均漏洞稳定在约 20 个——占比下降是结构改善，不是漏洞变少（2022 原文）", 0),
    ("按 C/C++ 历史密度 >1 个/kLOC 推算，Rust 已为 Android **预防数百个漏洞**（2022/2024 官方）", 0),
], t=Inches(6.35), size=12.5, gap=4)
add_footer(s, "security.googleblog.com 2022-12、2024-09、2024-10、2025-11（均已回源）", 5)

# ============ S6 洞察3 互操作 ============
s = content_slide("洞察 ③ · 工程门槛", "真正的门槛是互操作——Google 把它从信仰问题变成工程问题", 6)
bullets(s, [
    ("**先量化、再投入**（2021-06 原文）：用 objdump 分析最常用 C++ 库的导出函数参数类型", 0),
    ("核心库 **81% 类型即可对接、87% 可低成本对接**；Mainline 模块 88%/90%——可行性先被证明，才有大规模投入", 1),
    ("**工具链打底**：bindgen / cxx / autocxx / AIDL Rust 后端；粗粒度 FFI 哲学——锁、句柄等状态不跨语言", 0),
    ("**生态投资**：向 Rust Foundation 捐资 **100 万美元**设 Interop Initiative（2024-02）；工具覆盖 cbindgen/diplomat/crubit", 0),
    ("**依赖治理**：454 个第三方 crates 统一 vendored 进 `android-crates-io`（359 万行），聚合安全审计并公开（2024-02）", 0),
    ("**观点：渐进改造大型 C/C++ 系统，“与存量共存的能力”比“目标语言的优劣”更决定成败**", 0),
], t=Inches(2.0), size=14.5)
add_footer(s, "security.googleblog.com 2021-06、2024-02（均已回源）；AOSP main 实测", 6)

# ============ S7 洞察4 收益外溢 ============
s = content_slide("洞察 ④ · 采用动力", "收益已外溢到效率与资源：安全语言开始反哺交付", 7)
bullets(s, [
    ("交付效率（2025-11 官方 DORA 口径，同期同规模变更对比）", 0),
    ("代码评审耗时 **-25%**；返工轮次 **-20%**；中大型变更回滚率 **低 4 倍**", 1),
    ("官方定性：“the safer path is now also the faster one.”", 1),
    ("资源红利（2022-12 原文）", 0),
    ("UWB 新栈无需独立隔离进程：**省数 MB 内存、免一次 IPC 往返**；DoH3 用 async/await 单线程安全跑多任务，**线程数更少**", 1),
    ("增量拐点：第一方代码**净增行数 Rust 已超过 C++**（2025-11 官方图表）", 0),
    ("**观点：当安全带来的不是成本而是效率，采用就从“合规驱动”变成“利益驱动”——这是最强的自我加速信号**", 0),
], t=Inches(2.0), size=14.5)
add_footer(s, "security.googleblog.com 2022-12、2025-11《move fast and fix things》（均已回源）", 7)

# ============ S8 洞察5 路径 ============
s = content_slide("洞察 ⑤ · 落地路径", "七年路线：风险从高到低，环境从易到难", 8)
rows = [
    ("2019", "定向", "内部决策：新开发转向内存安全语言；18 个月搭工具链/测试/培训"),
    ("2021", "系统服务", "Android 12：平台级 Rust 支持发布；**Keystore2** 上线——首个旗舰 Rust 组件"),
    ("2022", "协议栈/虚拟化", "Android 13：UWB、DoH3（quiche）、AVF；150 万行、新增原生代码 21%"),
    ("2023", "裸机固件", "pvmfw 从 U-Boot（C）整体迁到 Rust——进入无 OS/无标准库环境"),
    ("2025", "内核", "6.12 GKI 内核：首个启用 Rust 支持、首个量产 Rust 驱动；GPU 驱动 Tyr 推进中"),
    ("2026", "基带固件", "Pixel 10 基带集成 Rust DNS 解析器（hickory-proto）——可远程攻击面的最深处"),
]
t_y = Inches(2.0)
for i, (yr, layer, desc) in enumerate(rows):
    y = t_y + Inches(0.82) * i
    tb = textbox(s, Inches(0.6), y, Inches(0.95), Inches(0.5))
    _fill_para(tb.text_frame.paragraphs[0], yr, 15, ACCENT, bold_all=True)
    rect(s, Inches(1.62), y + Inches(0.06), Inches(1.55), Inches(0.62), LIGHT)
    tb = textbox(s, Inches(1.66), y + Inches(0.14), Inches(1.5), Inches(0.5))
    _fill_para(tb.text_frame.paragraphs[0], layer, 11.5, ACCENT2, bold_all=True)
    tb = textbox(s, Inches(3.35), y + Inches(0.05), Inches(9.6), Inches(0.75))
    _fill_para(tb.text_frame.paragraphs[0], desc, 13, BODY)
bullets(s, [
    ("每一步都踩在“**解析不可信输入 / 权限最高**”的位置；越底层越晚——no_std、无分配器、厂商代码耦合，需要专门方法论（固件渐进替换教程、上游化 no_std）", 0),
], t=Inches(6.85), size=12.5, gap=2)
add_footer(s, "security.googleblog.com 2021–2026 各篇（均已回源）；source.android.com", 8)

# ============ S9 量与位置（图） ============
s = content_slide("规模 · 量与位置", "Rust 已是 AOSP 的结构性存在，但 C/C++ 仍是主体", 9)
chart_data = CategoryChartData()
chart_data.categories = ["2021-06\n官方", "2022-12\n官方", "2025-11\n官方", "2026-06\n全树统计"]
chart_data.add_series("Rust 代码量（百万行）", (0.1, 1.5, 5.0, 6.95))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(2.0), Inches(5.9), Inches(4.15), chart_data)
ch = gf.chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = "Rust 总量：10万 → 150万 → 500万 → 695万行"
for r in ch.chart_title.text_frame.paragraphs[0].runs:
    _set_font(r, 13, True, INK)
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.00"M"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(11)
plot.data_labels.font.bold = True
ser = ch.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = ACCENT
tb = textbox(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(4.3))
tf = tb.text_frame
lines = [
    ("**存量对比（第三方全树统计）**", 13.5, INK),
    ("C/C++ 占全树代码 54.7% → 49.0%（AOSP 14→17），占比下降但绝对量仍涨——存量按策略继续老化", 12, BODY),
    ("**第一方模块实测（AOSP main，2026-08）**", 13.5, INK),
    ("crosvm 38.9万 · 蓝牙 7.4万 · AVF 6.0万 · Nearby 5.7万 · keystore2 5.1万 · KeyMint 库 2.3万 · UWB 2.8万 · binder_rs 1.1万（采样合计约 76 万行）", 12, BODY),
    ("**依赖生态**：454 个 vendored crates = 359 万行，统一审计", 12, BODY),
    ("**观点：Rust 没有消灭 C++，它接管了 C++ 增量中风险最高的那部分**", 13.5, ACCENT),
]
first = True
for txt_, sz, col in lines:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_after = Pt(9)
    p.line_spacing = 1.15
    _fill_para(p, txt_, sz, col)
add_footer(s, "官方口径（2021-06/2022-12/2025-11 原文）；derdilla 全树统计；本报告克隆实测", 9)

# ============ S10 场景×特性 ============
s = content_slide("机制 · 场景 × 特性", "先落地的都是同一类位置；吃到的红利各不相同", 10)
tbl_rows = [
    ("场景", "代表模块（实测行数）", "吃到的 Rust 特性红利"),
    ("解析不可信输入", "UWB 2.8万 · DoH3 · crabbyavif 2.1万 · 基带 DNS", "边界检查消除整类漏洞；Result 强制错误处理；unsafe 收拢（UWB 全栈仅 2 处）"),
    ("密钥与 TEE", "keystore2 5.1万 · KeyMint 2.3万 · SecretKeeper", "所有权建模密钥生命周期；Send/Sync 并发安全；no_std 写到 TEE"),
    ("虚拟化与固件", "crosvm 38.9万 · AVF 6.0万 · pvmfw", "unsafe 边界可审计化；零成本抽象；no_std 进固件"),
    ("IPC 与互操作", "binder_rs 1.1万 · AIDL Rust 后端 · cxx", "零成本 FFI；接口契约编译期一致；tokio 异步"),
    ("内核与基带", "6.12 内核 Rust 驱动 · Pixel 10 基带", "类型编码内核/硬件 API 契约；无 GC 延迟可预测"),
]
from pptx.util import Cm
tbl_shape = s.shapes.add_table(len(tbl_rows), 3, Inches(0.6), Inches(2.0), Inches(12.15), Inches(4.3))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.5)
tbl.columns[1].width = Inches(4.35)
tbl.columns[2].width = Inches(5.3)
for ri, row in enumerate(tbl_rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.margin_top = cell.margin_bottom = Pt(4)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _fill_para(p, val, 12 if ri else 12.5, WHITE if ri == 0 else BODY, bold_all=(ri == 0))
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT2 if ri == 0 else (LIGHT if ri % 2 else WHITE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
bullets(s, [
    ("官方七特性清单（2021-04）：内存安全 / 并发安全 / 类型表达力 / 默认不可变 / 强制错误处理 / 强制初始化 / 整数安全——**Rust 的卖点从来不是一个词**", 0),
], t=Inches(6.5), size=12.5, gap=2)
add_footer(s, "security.googleblog.com 2021-04/2022-12/2023-10（均已回源）；本报告实测", 10)

# ============ S11 边界 ============
s = content_slide("边界 · 清醒认知", "Rust 不是银弹：价值在于把审计面缩小两个数量级", 11)
bullets(s, [
    ("**4%** 的 Rust 代码在 `unsafe{}` 块内（2025-11 官方）——这是全部“可能不安全”的表面积", 0),
    ("未遂事件（2025-11 官方披露）：CrabbyAVIF 一处 unsafe 代码存在线性缓冲区溢出", 0),
    ("发布前被拦截；Android 默认的 **Scudo 硬化分配器**确定性使其不可利用——纵深防御依然必要", 1),
    ("官方叙事随之升级：从“零漏洞”改为“**漏洞密度低 1000 倍**”——更诚实，也更可持续", 1),
    ("工程代价是真实的：基带案例需要给上游 crate 补 no_std、解决弱符号冲突、付出 371KB 体积", 0),
    ("**观点：不要用“绝对安全”推销 Rust——用“把安全审计面缩小到可管理范围”来定位它**", 0),
], t=Inches(2.0), size=14.5)
add_footer(s, "security.googleblog.com 2025-11、2026-04（均已回源）", 11)

# ============ S12 启示 ============
s = content_slide("启示 · 对从业者", "一套可复制的动作清单", 12)
acts = [
    ("1  先量化，后站队", "对自有 C/C++ 接口做 objdump 式类型分析，量化可互操作比例，用数据决定改造范围"),
    ("2  增量默认安全语言", "新功能/新模块默认内存安全语言；重写只碰“新 + 高风险 + 边界清晰”的存量（如解析器、固件）"),
    ("3  治理依赖与审计", "第三方 crates 统一 vendored、聚合审计并公开；unsafe 代码收拢到最小可审面"),
    ("4  配套工程体系", "构建系统一等支持、CI 强制 fmt/clippy、培训课程化（Google 500+ 工程师、96% 好评）"),
    ("5  用效率指标闭环", "用 DORA 指标（评审时长、返工、回滚率）证明“更安全 = 更快”，让采用从合规变利益"),
]
top = Inches(2.05)
for i, (h, b) in enumerate(acts):
    y = top + Inches(0.98) * i
    rect(s, Inches(0.6), y, Inches(12.15), Inches(0.86), LIGHT)
    tb = textbox(s, Inches(0.85), y + Inches(0.05), Inches(3.3), Inches(0.7))
    _fill_para(tb.text_frame.paragraphs[0], h, 14, ACCENT, bold_all=True)
    tb = textbox(s, Inches(4.3), y + Inches(0.08), Inches(8.4), Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    _fill_para(tf.paragraphs[0], b, 12, BODY)
add_footer(s, "综合官方 13 篇原文与本报告分析", 12)

# ============ S13 结语 ============
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, Inches(0.16), SH, ACCENT)
tb = textbox(s, Inches(0.9), Inches(2.3), Inches(11.6), Inches(2.4))
tf = tb.text_frame
tf.word_wrap = True
_fill_para(tf.paragraphs[0], "大型 C/C++ 系统的内存安全问题是可以被“设计掉”的——", 26, WHITE, bold_all=True)
p = tf.add_paragraph()
_fill_para(p, "前提是你接受渐进主义。", 26, RGBColor(0xF0, 0x8A, 0x4B), bold_all=True)
tb = textbox(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.4))
tf = tb.text_frame
_fill_para(tf.paragraphs[0],
           "漏洞会衰老，增量会置换，七年之后曲线自然拐头。",
           15, RGBColor(0xC9, 0xD2, 0xDA))
p = tf.add_paragraph()
_fill_para(p, "Rust in Android 2019–2026 · 官方 13 篇原文回源 · AOSP main 实测 · 全树统计交叉验证",
           11, RGBColor(0x9A, 0xA5, 0xB0))

out = r"F:\lee_space\code\research\rust-aosp-insights.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
