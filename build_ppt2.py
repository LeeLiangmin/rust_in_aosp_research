# -*- coding: utf-8 -*-
"""Insight deck v2, structured after rust-in-aosp.md."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

INK = RGBColor(0x1F, 0x2A, 0x37)
BODY = RGBColor(0x37, 0x42, 0x4F)
MUTED = RGBColor(0x6B, 0x75, 0x80)
ACCENT = RGBColor(0xB7, 0x41, 0x0E)
ACCENT2 = RGBColor(0x14, 0x5A, 0x6E)
LIGHT = RGBColor(0xF4, 0xF1, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set_font(run, size=14, bold=False, color=BODY, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", FONT)


def _fill_para(p, text, size, color, bold_all=False):
    for i, seg in enumerate(text.split("**")):
        if not seg:
            continue
        r = p.add_run()
        r.text = seg
        _set_font(r, size=size, bold=bold_all or (i % 2 == 1), color=color)


def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def rect(slide, l, t, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def add_footer(slide, note, page):
    tb = textbox(slide, Inches(0.55), Inches(7.05), Inches(11.4), Inches(0.38))
    _fill_para(tb.text_frame.paragraphs[0], "数据来源：" + note, 9, MUTED)
    tb2 = textbox(slide, Inches(12.35), Inches(7.05), Inches(0.7), Inches(0.38))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    _fill_para(p2, str(len(prs.slides._sldIdLst)), 10, MUTED, bold_all=True)


def content_slide(kicker, title, page=None, title_size=25):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, Inches(0.16), SH, ACCENT)
    tb = textbox(s, Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.42))
    _fill_para(tb.text_frame.paragraphs[0], kicker, 12.5, ACCENT, bold_all=True)
    tb = textbox(s, Inches(0.55), Inches(0.68), Inches(12.4), Inches(1.1))
    _fill_para(tb.text_frame.paragraphs[0], title, title_size, INK, bold_all=True)
    rect(s, Inches(0.57), Inches(1.72), Inches(1.35), Pt(3.2), ACCENT)
    return s


def bullets(slide, items, l=Inches(0.6), t=Inches(1.95), w=Inches(12.15), h=Inches(4.9), size=14.5, gap=9):
    tb = textbox(slide, l, t, w, h)
    tf = tb.text_frame
    first = True
    for text, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap if lvl == 0 else 4)
        p.line_spacing = 1.16
        mark = "▍ " if lvl == 0 else "– "
        r = p.add_run()
        r.text = "    " * lvl + mark
        _set_font(r, size=size if lvl == 0 else size - 1.5, bold=(lvl == 0),
                  color=ACCENT if lvl == 0 else MUTED)
        _fill_para(p, text, size if lvl == 0 else size - 1.5, BODY if lvl else INK)
    return tb


def make_table(slide, rows, l, t, w, h, widths, header_fill=ACCENT2, size=11.5):
    shp = slide.shapes.add_table(len(rows), len(rows[0]), l, t, w, h)
    tbl = shp.table
    for ci, cw in enumerate(widths):
        tbl.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_top = cell.margin_bottom = Pt(3)
            tf = cell.text_frame
            tf.word_wrap = True
            _fill_para(tf.paragraphs[0], val, size + 0.5 if ri == 0 else size,
                       WHITE if ri == 0 else BODY, bold_all=(ri == 0))
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if ri == 0 else (LIGHT if ri % 2 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl


# ============ S1 封面 ============
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(6.9), SW, Inches(0.6), ACCENT)
tb = textbox(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.5))
_fill_para(tb.text_frame.paragraphs[0], "洞察报告 · 基于《Rust 在 Android/AOSP 中的调研报告》· 2026-08", 13, RGBColor(0xE8, 0xB4, 0x9B), bold_all=True)
tb = textbox(s, Inches(0.9), Inches(1.85), Inches(11.6), Inches(2.6))
tf = tb.text_frame
_fill_para(tf.paragraphs[0], "Rust 在 Android 的七年：", 40, WHITE, bold_all=True)
p = tf.add_paragraph()
_fill_para(p, "一场方法论的胜利", 40, RGBColor(0xF0, 0x8A, 0x4B), bold_all=True)
tb = textbox(s, Inches(0.9), Inches(4.3), Inches(11.4), Inches(1.3))
tf = tb.text_frame
_fill_para(tf.paragraphs[0], "内存安全漏洞占比 76%（2019）→ 35%（2022）→ 24%（2024）→ <20%（2025）", 17, RGBColor(0xC9, 0xD2, 0xDA))
p = tf.add_paragraph()
_fill_para(p, "改变结果的不是一门语言，而是一套可复制的渐进改造方法。", 17, RGBColor(0xC9, 0xD2, 0xDA))
tb = textbox(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5))
_fill_para(tb.text_frame.paragraphs[0], "数据底座：Google 官方博客 13 篇逐句回源 · AOSP main 45+ 仓库克隆实测 · 第三方全树统计（AOSP 14/16/17）", 11, RGBColor(0x9A, 0xA5, 0xB0))

# ============ S2 摘要 ============
s = content_slide("摘要 · 报告核心结论", "Rust 是 Google 内存安全战略在原生层的载体：新代码置换，而非存量重写", 2)
bullets(s, [
    ("动因：2019 年内存类缺陷占 Android 漏洞 **90%**（追溯口径 76%）；缓解手段（ASLR/sanitizer/沙箱）只抬高利用成本，无法消除漏洞", 0),
    ("策略：**新代码用内存安全语言 + 不重写成熟存量 + 互操作优先 + 风险驱动选点**，配合内置工具链与统一依赖治理", 0),
    ("现状（2026）：Rust 覆盖工具链、IPC、密钥/TEE、虚拟化、协议栈、媒体、内核乃至**基带固件**；官方总量约 **500 万行**，全树统计 **695 万行（3.8%）**，本报告实测第一方约 **76 万行**", 0),
    ("效果：漏洞占比 76%→**<20%**、绝对数量 223→年化 36；Rust 漏洞密度比 C/C++ **低三个数量级**；评审快 25%、回滚低 4 倍", 0),
    ("**一句话：Android 用七年证明，巨型 C/C++ 系统的内存安全问题可以被“设计掉”——前提是接受渐进主义**", 0),
], t=Inches(2.0), size=14.5, gap=11)
add_footer(s, "rust-in-aosp.md 摘要；官方 13 篇原文回源核对", 2)

# ============ S3 动因 ============
s = content_slide("一、落地过程 · 动因", "内存安全是 Android 漏洞的结构性来源，而缓解手段已到天花板", 3)
bullets(s, [
    ("2019 年官方盘点（原文已核对）：UAF、整数溢出、越界读写合计占当年漏洞 **90%**，数组边界类占 34%", 0),
    ("“76% 为内存安全漏洞”是 Google 2022/2024 回顾的追溯口径；2021 年 Rust 公告用的是“约 70% 高危漏洞”", 1),
    ("为什么前两代手段不够", 0),
    ("被动修补＝打地鼠，修复本身还引入新漏洞；运行时缓解＝军备竞赛，每代缓解都有对应绕过", 1),
    ("沙箱（Rule of 2）吃内存、加 IPC 延迟，且不消除漏洞——密度高时可被串联打穿（2021 原文）", 1),
    ("为什么是 Rust 而不是别的：Java/Kotlin 内存安全但有 GC，管不到 HAL/守护进程/固件/内核；Go 同样受 GC 与运行时拖累", 0),
    ("**Rust 是唯一同时满足“内存安全 + 无 GC + 零成本抽象 + C 级性能 + 成熟生态”的系统语言**", 1),
], t=Inches(2.0), size=14.5, gap=9)
add_footer(s, "security.googleblog.com 2019-05、2021-04（均已回源）；rust-in-aosp.md §1.1", 3)

# ============ S4 策略 ============
s = content_slide("一、落地过程 · 策略", "四条策略：漏洞会“衰老”，所以只换增量、不动存量", 4)
bullets(s, [
    ("① **新代码优先，不重写存量**——依据是漏洞半衰期：约 50% 内存 bug 诞生于不足一年的新代码；5 年陈代码漏洞密度低 **3.4–7.4 倍**", 0),
    ("② **互操作优先，而非替换**——先量化再投入：核心库 81%/87%、Mainline 88%/90% 类型可对接；bindgen/cxx/AIDL 后端打底；2024 年再捐 **100 万美元**攻关", 0),
    ("③ **风险驱动选点**——先打攻击面最大的位置：解析不可信输入的组件、特权服务、隔离边界（TEE/固件/虚拟化）", 0),
    ("④ **工具链与流程配套**——prebuilts/rust 工具链、Soong 原生 rust_* 模块、CI 强制 rustfmt/clippy、454 个第三方 crates 统一 vendored 与聚合审计", 0),
    ("**观点：安全改造的关键决策不是“用什么语言”，而是“改哪部分代码”**", 0),
], t=Inches(2.0), size=14.5, gap=10)
add_footer(s, "security.googleblog.com 2021-04、2021-06、2024-02、2024-09（均已回源）；rust-in-aosp.md §1.2", 4)

# ============ S5 时间线 ============
s = content_slide("一、落地过程 · 时间线", "七年路线：系统服务 → 协议栈 → 虚拟化 → 固件 → 内核 → 基带", 5)
rows = [
    ("2019", "定向", "内存类缺陷占漏洞 90%；团队决定新开发转向内存安全语言"),
    ("2021", "平台支持", "Android 12：Rust 平台级支持发布；Keystore2 上线；AOSP 内 Rust 超 10 万行"),
    ("2022", "规模落地", "Android 13：150 万行、新增原生代码 21% 为 Rust；UWB/DoH3/AVF；漏洞占比降至 35%"),
    ("2023", "TEE/裸机", "pvmfw 从 U-Boot（C）迁到 Rust；Trusty 可信应用 Rust 化；500+ 工程师受训"),
    ("2024", "方法论定型", "占比降至 24%（绝对数年化 36）；safe coding 体系化；固件渐进替换教程公开"),
    ("2025", "内核", "占比首次 <20%、总量约 500 万行、密度低 1000 倍；6.12 内核首个量产 Rust 驱动"),
    ("2026", "基带", "Pixel 10 基带集成 Rust DNS 解析器——Rust 进入可远程攻击面的最深处"),
]
t_y = Inches(1.95)
for i, (yr, layer, desc) in enumerate(rows):
    y = t_y + Inches(0.72) * i
    tb = textbox(s, Inches(0.6), y, Inches(0.9), Inches(0.45))
    _fill_para(tb.text_frame.paragraphs[0], yr, 14, ACCENT, bold_all=True)
    rect(s, Inches(1.58), y + Inches(0.02), Inches(1.45), Inches(0.56), LIGHT)
    tb = textbox(s, Inches(1.63), y + Inches(0.11), Inches(1.4), Inches(0.45))
    _fill_para(tb.text_frame.paragraphs[0], layer, 10.5, ACCENT2, bold_all=True)
    tb = textbox(s, Inches(3.2), y + Inches(0.02), Inches(9.75), Inches(0.68))
    _fill_para(tb.text_frame.paragraphs[0], desc, 12, BODY)
add_footer(s, "security.googleblog.com 2019–2026 各篇（13 篇已回源）；rust-in-aosp.md §1.3", 5)

# ============ S5b 演进位置（架构图） ============
s = content_slide("一、落地过程 · 演进位置", "Rust 的边界，就是 Java/Kotlin 管不到的地方")
img = r"F:\lee_space\code\research\综合\rust在aosp的落地分析\rust_in_aosp_architecture.png"
s.shapes.add_picture(img, Inches(0.6), Inches(1.9), width=Inches(4.3))
bullets(s, [
    ("**上三层（App / Framework / ART）未动**：Java/Kotlin 本就内存安全——不是遗漏，是策略（2021 原文）", 0),
    ("主战场在 ART 以下，自顶向下推进", 0),
    ("System services（2021）：Keystore2、KeyMint 等守护进程", 1),
    ("Native daemons & libraries（2021–22）：蓝牙、UWB、DNS-over-HTTP3", 1),
    ("HAL（2022+）：userspace HAL 迁移中（蓝牙 offload HAL 等）", 1),
    ("Kernel（2025）：6.12 GKI 首个量产 Rust 驱动、GPU 驱动 Tyr 推进中", 1),
    ("两块在官方架构图**之外**：AVF（跨 HAL/内核的独立子系统，2022–23）；固件/基带（OS 之外，2023–26）", 0),
    ("演进规律：每下一层都要先解决该层工程约束（no_std、代码体积、厂商构建系统）——**所以越底层越晚**", 0),
], l=Inches(5.2), t=Inches(1.95), w=Inches(7.55), size=13.5, gap=8)
add_footer(s, "图据 source.android.com 官方架构重绘并标注年份；rust-in-aosp.md §一/§三", 6)

# ============ S6 总量与口径 ============
s = content_slide("二、量与位置 · 总量", "三条数据链交叉验证：Rust 已是 AOSP 的结构性存在", 6)
chart_data = CategoryChartData()
chart_data.categories = ["2021-06", "2022-12", "2025-11", "2026-06\n(全树)"]
chart_data.add_series("Rust 代码量（百万行）", (0.1, 1.5, 5.0, 6.95))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.95), Inches(5.9), Inches(4.4), chart_data)
ch = gf.chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = "Rust 总量：10万→150万→500万→695万行"
for r in ch.chart_title.text_frame.paragraphs[0].runs:
    _set_font(r, 13, True, INK)
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.00"M"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(11)
plot.data_labels.font.bold = True
ser = ch.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = ACCENT
tb = textbox(s, Inches(6.85), Inches(1.95), Inches(6.05), Inches(4.6))
tf = tb.text_frame
lines = [
    ("**口径一 · 官方**（均已回源）：2021-06 超 10 万行 → Android 13 约 150 万行（含依赖）→ 2025-11 约 500 万行", 12.5, BODY),
    ("**口径二 · 全树统计**（derdilla，tokei）：AOSP 17 全树 Rust 695 万行，占全部代码约 3.8%", 12.5, BODY),
    ("**口径三 · 本报告实测**（main 分支克隆）：第一方平台约 76.2 万行；vendored 第三方 crates 454 个、359 万行", 12.5, BODY),
    ("**注意**：官方从未发布逐模块行数清单；逐模块规模来自本报告实测与社区清单（awesome-aosp-rust）", 11.5, MUTED),
    ("**观点：三个口径讲的是同一句话——Rust 在 AOSP 已不是试验，而是基础设施**", 13.5, ACCENT),
]
first = True
for txt_, sz, col in lines:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_after = Pt(10)
    p.line_spacing = 1.18
    _fill_para(p, txt_, sz, col)
add_footer(s, "官方 2021-06/2022-12/2025-11 原文；derdilla.com；G:\\aosp-research 实测；rust-in-aosp.md §2.1", 6)

# ============ S7 全树跨版本 ============
s = content_slide("二、量与位置 · 跨版本对比", "C/C++ 占比持续下降，Rust 在平台核心两年涨 3.6 倍", 7)
make_table(s, [
    ("版本（统计时间）", "Rust 总行数", "其中 Core 类 Rust", "Rust 占比", "C/C++ 行数", "C/C++ 占比", "全树总量", "设备上运行"),
    ("AOSP 14（2024-08）", "约 498 万", "7.8 万（0.72%）", "4.0%", "约 6836 万", "54.7%", "约 1.25 亿", "约 6600 万"),
    ("AOSP 16（2025-06）", "约 593 万", "13.6 万（1.25%）", "3.9%", "约 7864 万", "51.9%", "约 1.52 亿", "约 9200 万"),
    ("AOSP 17（2026-06）", "约 695 万", "35.8 万（2.70%）", "3.8%", "约 9009 万", "49.0%", "约 1.84 亿", "约 1.13 亿"),
], Inches(0.6), Inches(2.0), Inches(12.15), Inches(1.9),
    [Inches(2.1), Inches(1.35), Inches(1.85), Inches(1.1), Inches(1.5), Inches(1.25), Inches(1.5), Inches(1.5)], size=11.5)
bullets(s, [
    ("**Core 类 Rust 两年约 3.6 倍**（7.8 万→35.8 万）、Userspace 翻倍——这是第一方 Rust 化的真实信号", 0),
    ("Rust 全树占比略降（4.0%→3.8%）不是放缓：是 C/C++ 侧**测试代码**增长更快（原作者解释）", 0),
    ("C/C++ 占比持续下降（54.7%→49.0%）但绝对量仍增长——**存量继续存在并老化，与“不重写”策略一致**", 0),
    ("Third-party 类 Rust（389 万→483 万）主要反映 vendored crates 依赖膨胀", 0),
], t=Inches(4.35), size=13.5, gap=8)
add_footer(s, "derdilla.com AOSP 14/16/17 全树统计（tokei，无 Android 15 专篇）；rust-in-aosp.md §2.2", 7)

# ============ S8 模块地图 ============
s = content_slide("二、量与位置 · 模块地图", "Rust 遍布 AOSP 各层：从工具链到基带固件", 8)
make_table(s, [
    ("层次", "代表组件（实测行数，AOSP main 2026-08）"),
    ("基础设施/工具链", "prebuilts/rust 工具链 · Soong rust_* 模块 · AIDL Rust 后端 3.4万 · cxx 2.1万 · vendored crates 359万/454个"),
    ("IPC 与系统服务", "keystore2 5.1万（A12 起替换旧 keystore）· binder_rs/rpcbinder 1.1万 · librustutils/prefetch/debuggerd 绑定"),
    ("安全 / TEE / 固件", "KeyMint 支持库 2.3万 · Trusty 可信应用（KeyMint/SecretKeeper/AuthMgr）· libavb 0.5万"),
    ("虚拟化", "crosvm 38.9万（最大第一方单体）· AVF/Microdroid/pvmfw 6.0万"),
    ("连接与协议栈", "蓝牙 Gabeldorsche/offload/floss 7.4万 · Nearby beto-rust 5.7万 · UWB 2.8万 · DoH3（quiche）0.4万"),
    ("媒体与图形", "crabbyavif AVIF 解码 2.1万 · pica 矢量渲染 0.4万"),
    ("内核", "Linux 6.12 GKI：首个启用 Rust 支持、首个量产 Rust 驱动（GPU 驱动 Tyr 推进中）"),
], Inches(0.6), Inches(1.95), Inches(12.15), Inches(4.55),
    [Inches(2.5), Inches(9.65)], size=11.5)
add_footer(s, "本报告克隆实测（G:\\aosp-research / G:\\aosp-scan）；rust-in-aosp.md §2.3", 8)

# ============ S9 与 C/C++ 对比 ============
s = content_slide("二、量与位置 · 与 C/C++ 对比", "新旧分层：Rust 接管了 C++ 增量中风险最高的部分", 9)
bullets(s, [
    ("同一仓库内的新旧分层最能说明策略", 0),
    ("DnsResolver：新功能 DoH3 用 Rust（3,879 行），存量 resolver 仍是 C++（40,904 行）", 1),
    ("system/security：新 keystore2 用 Rust（50,865 行），已反超同仓库 C++（16,631 行）", 1),
    ("Trusty：老 storage/gatekeeper 仍是 C/C++（约 3 万行），新 KeyMint/SecretKeeper 是 Rust", 1),
    ("增量：A13 新增原生代码 **21%** 为 Rust；2025 年第一方**净增行数 Rust 已超过 C++**（官方图表）", 0),
    ("质量：漏洞占比 76%→35%→24%→**<20%**；绝对数量 **223→85→年化 36**；密度差 **1000 倍**（C/C++ 历史约 1000 个/MLOC）", 0),
    ("依赖生态：vendored crates 从 A13 时代数十万行膨胀到 **359 万行**——公共底座已经形成", 0),
], t=Inches(2.0), size=14.5, gap=9)
add_footer(s, "官方 2022-12/2025-11 原文；本报告实测；rust-in-aosp.md §2.4", 9)

# ============ S10 七特性 ============
s = content_slide("三、场景×特性 · 官方清单", "Rust 的卖点从来不是一个词：官方七特性清单", 10)
feats = [
    ("内存安全", "所有权/借用检查 + 少量运行时检查，整类消灭 UAF/越界/双重释放"),
    ("数据并发安全", "Send/Sync 编译期杜绝数据竞争——“Fearless Concurrency”"),
    ("类型表达力", "newtype / 带数据的枚举，从类型层面排除非法状态"),
    ("默认不可变", "最小权限原则；官方吐槽 C++ 的 const“用得少且不一致”"),
    ("强制错误处理", "Result + ?；反例：提权漏洞 Rage Against the Cage 根因即未检查错误返回值"),
    ("强制初始化", "未初始化内存历史上占 Android 漏洞 3–5%；A11 的自动初始化只是缓解"),
    ("整数安全", "转换必须显式 cast、无隐式截断；调试构建默认开启溢出检查"),
]
for i, (h, b) in enumerate(feats):
    y = Inches(1.95) + Inches(0.72) * i
    rect(s, Inches(0.6), y, Inches(2.6), Inches(0.62), LIGHT)
    tb = textbox(s, Inches(0.75), y + Inches(0.1), Inches(2.4), Inches(0.5))
    _fill_para(tb.text_frame.paragraphs[0], h, 13, ACCENT2, bold_all=True)
    tb = textbox(s, Inches(3.45), y + Inches(0.06), Inches(9.4), Inches(0.66))
    _fill_para(tb.text_frame.paragraphs[0], b, 12, BODY)
add_footer(s, "security.googleblog.com 2021-04 “Prioritizing prevention”（已回源）；rust-in-aosp.md §三", 10)

# ============ S11 场景映射 ============
s = content_slide("三、场景×特性 · 六大场景", "先落地的都是同一类位置，但吃到的红利各不相同", 11)
make_table(s, [
    ("场景", "代表模块", "痛点", "吃到的特性红利"),
    ("解析不可信输入", "UWB/DoH3/AVIF/AVB/基带DNS", "C/C++ 解析器是漏洞头号产地", "边界检查整类消除；Result 强制错误处理；UWB 全栈仅 2 处 unsafe"),
    ("密钥与 TEE", "keystore2/KeyMint/SecretKeeper", "最高价值资产 + 并发服务", "所有权建模密钥生命周期；Send/Sync；no_std 写到 TEE"),
    ("虚拟化与固件", "crosvm/AVF/pvmfw", "VM 逃逸攻击面；固件信任根", "unsafe 收拢可审计；零成本抽象无 GC；no_std（U-Boot 反面教材）"),
    ("IPC 与互操作", "binder_rs/AIDL/cxx", "新旧代码必须共存", "零成本 FFI；接口契约编译期一致；tokio 异步"),
    ("内核与基带", "6.12 Rust 驱动/Pixel 10 基带", "最高权限环境，出错即沦陷", "类型编码内核 API 契约；无 GC 延迟可预测"),
    ("性能敏感新组件", "UWB 时序/pica/libprefetch", "不能拿性能换安全", "零成本抽象、无 GC——不必做这笔交易"),
], Inches(0.6), Inches(1.95), Inches(12.15), Inches(4.6),
    [Inches(2.2), Inches(3.0), Inches(3.0), Inches(3.95)], size=10.5)
add_footer(s, "rust-in-aosp.md §三 场景 1–6；官方 2022-12/2023-10/2026-04 原文", 11)

# ============ S12 实证结论 ============
s = content_slide("三、实证结论", "收益不止安全：资源、效率、审计面三重红利", 12)
bullets(s, [
    ("**资源红利**（2022-12 原文）：UWB 新栈无需独立隔离进程——省数 MB 内存、免一次 IPC 往返；DoH3 用 async/await 单线程安全跑多任务，线程更少", 0),
    ("**交付效率**（2025-11 官方 DORA 口径）：评审耗时 **-25%**、返工轮次 **-20%**、中大型变更回滚率 **低 4 倍**——“the safer path is now also the faster one”", 0),
    ("**审计面收拢**：全平台仅约 **4%** Rust 代码在 unsafe 块内；CrabbyAVIF 未遂事件发布前被拦截、且被 Scudo 硬化分配器确定性中和", 0),
    ("**官方叙事升级**：从“零漏洞”改为“漏洞密度低 1000 倍”——更诚实，也更可持续", 0),
    ("**观点：当安全带来的不是成本而是效率，采用就从合规驱动变成利益驱动——这是最强的自我加速信号**", 0),
], t=Inches(2.0), size=14.5, gap=11)
add_footer(s, "security.googleblog.com 2022-12、2025-11（均已回源）；rust-in-aosp.md §三 实证结论", 12)

# ============ S13 边界与数据链 ============
s = content_slide("边界 · 数据来源", "不是银弹；三条数据链支撑每一个数字", 13)
bullets(s, [
    ("边界认知", 0),
    ("Rust 的价值不是“绝对安全”，而是**把安全审计面缩小到可管理范围**（4% unsafe + 纵深防御）", 1),
    ("工程代价真实存在：基带案例需给上游 crate 补 no_std、解决弱符号冲突、付出 371KB 体积", 1),
    ("三条互相印证的数据链", 0),
    ("官方原文 13 篇逐句回源（2019–2026）——趋势、口径、策略表述", 1),
    ("第三方全树统计（derdilla，AOSP 14/16/17，tokei）——跨版本语言占比", 1),
    ("本报告实测（main 分支 45+ 仓库克隆，G:\\aosp-research / G:\\aosp-scan）——模块级行数", 1),
    ("局限：采样非全量（NFC 等已验证为 0）；行数含注释/空行；vendored 依赖与第一方分开计", 0),
], t=Inches(2.0), size=14, gap=9)
add_footer(s, "rust-in-aosp.md §四", 13)

# ============ S14 结语 ============
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, Inches(0.16), SH, ACCENT)
tb = textbox(s, Inches(0.9), Inches(2.15), Inches(11.6), Inches(2.6))
tf = tb.text_frame
tf.word_wrap = True
_fill_para(tf.paragraphs[0], "大型 C/C++ 系统的内存安全问题可以被“设计掉”——", 26, WHITE, bold_all=True)
p = tf.add_paragraph()
_fill_para(p, "前提是你接受渐进主义。", 26, RGBColor(0xF0, 0x8A, 0x4B), bold_all=True)
tb = textbox(s, Inches(0.9), Inches(4.65), Inches(11.5), Inches(1.5))
tf = tb.text_frame
_fill_para(tf.paragraphs[0], "漏洞会衰老，增量会置换，七年之后曲线自然拐头。", 15, RGBColor(0xC9, 0xD2, 0xDA))
p = tf.add_paragraph()
_fill_para(p, "详见完整调研报告：rust-in-aosp.md（落地过程 · 量与位置 · 场景×特性 · 数据来源）", 12, RGBColor(0x9A, 0xA5, 0xB0))

out = r"F:\lee_space\code\research\rust-aosp-insights.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
