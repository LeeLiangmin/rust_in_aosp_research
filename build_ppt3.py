# -*- coding: utf-8 -*-
"""Insight deck v3: 5 pages per user spec.
1. AOSP 架构（ART 之下用 Rust） 2. 演进：位置越来越底层、围栏越来越少
3. 代码量绝对趋势 + 对比 C++ 相对趋势  4. 混合编程方案  5. 实际项目 × Rust 特性
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

INK = RGBColor(0x1F, 0x2A, 0x37)
BODY = RGBColor(0x37, 0x42, 0x4F)
MUTED = RGBColor(0x6B, 0x75, 0x80)
ACCENT = RGBColor(0xB7, 0x41, 0x0E)      # rust orange
ACCENT2 = RGBColor(0x14, 0x5A, 0x6E)     # teal
ACCENT3 = RGBColor(0x3E, 0x6B, 0x34)     # green
LIGHT = RGBColor(0xF4, 0xF1, 0xEC)
LIGHT2 = RGBColor(0xE9, 0xF0, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set_font(run, size=14, bold=False, color=BODY, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", FONT)


def _fill_para(p, text, size, color, bold_all=False):
    for i, seg in enumerate(text.split("**")):
        if not seg:
            continue
        r = p.add_run()
        r.text = seg
        _set_font(r, size=size, bold=bold_all or (i % 2 == 1), color=color)


def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def rect(slide, l, t, w, h, fill, line=None, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    if rounded:
        try:
            sp.adjustments[0] = 0.06
        except Exception:
            pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def shape_text(sp, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.05
        _fill_para(p, ln, size, color, bold_all=bold)


def add_footer(slide, note, page):
    tb = textbox(slide, Inches(0.55), Inches(7.08), Inches(11.4), Inches(0.36))
    _fill_para(tb.text_frame.paragraphs[0], "数据来源：" + note, 9, MUTED)
    tb2 = textbox(slide, Inches(12.45), Inches(7.08), Inches(0.6), Inches(0.36))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    _fill_para(p2, str(page), 10, MUTED, bold_all=True)


def content_slide(kicker, title, page, title_size=25):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, Inches(0.16), SH, ACCENT)
    tb = textbox(s, Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.42))
    _fill_para(tb.text_frame.paragraphs[0], kicker, 12.5, ACCENT, bold_all=True)
    tb = textbox(s, Inches(0.55), Inches(0.68), Inches(12.4), Inches(1.1))
    _fill_para(tb.text_frame.paragraphs[0], title, title_size, INK, bold_all=True)
    rect(s, Inches(0.57), Inches(1.7), Inches(1.35), Pt(3.2), ACCENT)
    return s


def bullets(slide, items, l=Inches(0.6), t=Inches(1.95), w=Inches(12.15), h=Inches(4.9), size=14.5, gap=9):
    tb = textbox(slide, l, t, w, h)
    tf = tb.text_frame
    first = True
    for text, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap if lvl == 0 else 4)
        p.line_spacing = 1.16
        mark = "▍ " if lvl == 0 else "– "
        r = p.add_run()
        r.text = "    " * lvl + mark
        _set_font(r, size=size if lvl == 0 else size - 1.5, bold=(lvl == 0),
                  color=ACCENT if lvl == 0 else MUTED)
        _fill_para(p, text, size if lvl == 0 else size - 1.5, BODY if lvl else INK)
    return tb


def make_table(slide, rows, l, t, w, h, widths, header_fill=ACCENT2, size=11.5):
    shp = slide.shapes.add_table(len(rows), len(rows[0]), l, t, w, h)
    tbl = shp.table
    for ci, cw in enumerate(widths):
        tbl.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_top = cell.margin_bottom = Pt(3)
            tf = cell.text_frame
            tf.word_wrap = True
            _fill_para(tf.paragraphs[0], val, size + 0.5 if ri == 0 else size,
                       WHITE if ri == 0 else BODY, bold_all=(ri == 0))
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if ri == 0 else (LIGHT if ri % 2 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl


# ============ S1 封面 ============
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(6.9), SW, Inches(0.6), ACCENT)
tb = textbox(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.5))
_fill_para(tb.text_frame.paragraphs[0], "调研报告 · 基于 AOSP main 源码统计与 Google 官方博客 13 篇原文 · 2026-08", 13, RGBColor(0xE8, 0xB4, 0x9B), bold_all=True)
tb = textbox(s, Inches(0.9), Inches(1.9), Inches(11.6), Inches(2.4))
tf = tb.text_frame
_fill_para(tf.paragraphs[0], "Rust 在 Android 中的落地情况", 40, WHITE, bold_all=True)
p = tf.add_paragraph()
_fill_para(p, "位置逐年下沉，规模持续增长", 40, RGBColor(0xF0, 0x8A, 0x4B), bold_all=True)
tb = textbox(s, Inches(0.9), Inches(4.35), Inches(11.6), Inches(1.6))
tf = tb.text_frame
first = True
for ln in [
    "内容：架构位置 → 七年演进 → 代码量趋势 → 混合编程方案 → 实际项目分析",
    "两个主要结论：Rust 进入的层越来越低、隔离措施逐步减少；收益不只是安全，还有开发效率",
]:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_after = Pt(8)
    p.line_spacing = 1.15
    _fill_para(p, ln, 16, RGBColor(0xC9, 0xD2, 0xDA))
tb = textbox(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5))
_fill_para(tb.text_frame.paragraphs[0], "数据基础：rust-in-aosp.md + 五个代表性库的源码分析（G:\\aosp-scan 45+ 仓库）", 11, RGBColor(0x9A, 0xA5, 0xB0))

# ============ S2 架构：ART 之下 ============
s = content_slide("第 1 页 · AOSP 架构", "Rust 用在 ART 以下的层：ART 之上 Java/Kotlin 已经是内存安全语言", 2)

# 左侧：六层架构图
LX, LW = Inches(0.6), Inches(7.6)
layers = [
    ("① 应用层", "Java / Kotlin", "本来就是内存安全语言", False),
    ("② 应用框架（system_server）", "Java / Kotlin + JNI", "本来就是内存安全语言", False),
    ("③ ART / 运行时", "Java 托管执行", "有 GC——Rust 不进入这一层及以上", False),
    ("④ Native 库与守护进程", "keystore2 · 蓝牙 gd · UWB · DoH3 · libbinder_rs · libprefetch", "Rust 代码量最大的一层（2021 起）", True),
    ("⑤ HAL（AIDL）", "KeyMint HAL · 蓝牙 offload HAL", "AIDL 接口不变，只换实现语言（2022 起）", True),
    ("⑥ 内核 / 固件 / TEE / 基带", "ashmem 驱动 · pvmfw · Trusty TA · libavb · Pixel 基带", "无操作系统环境，需要 no_std（2023–26）", True),
]
y = Inches(1.95)
hh = Inches(0.78)
for name, impl, note, is_rust in layers:
    band = rect(s, LX, y, LW, hh - Inches(0.06), ACCENT if is_rust else LIGHT, rounded=True)
    if is_rust:
        shape_text(band, "**" + name + "**\n" + impl, 10.5, WHITE, align=PP_ALIGN.LEFT)
    else:
        shape_text(band, "**" + name + "**    " + impl, 10.5, BODY, align=PP_ALIGN.LEFT)
    tag = rect(s, LX + LW + Inches(0.12), y + Inches(0.12), Inches(4.4), hh - Inches(0.3),
               LIGHT2 if is_rust else WHITE, line=None, rounded=True)
    shape_text(tag, note, 10.5, ACCENT2 if is_rust else MUTED, bold=is_rust)
    y += hh

tb = textbox(s, Inches(0.6), y + Inches(0.06), Inches(12.2), Inches(0.7))
_fill_para(tb.text_frame.paragraphs[0],
           "原则：只在 Java/Kotlin 覆盖不到的地方用 Rust。③层及以上已经有 GC 语言保证内存安全，没有替换的必要（2021 官方原文）",
           12.5, INK, bold_all=False)
add_footer(s, "rust-in-aosp.md §一；source.android.com 架构文档；security.googleblog.com 2021-04", 2)

# ============ S3 演进：越来越底层、围栏越来越少 ============
s = content_slide("第 2 页 · 七年演进", "进入的层越来越低，配套的隔离措施越来越少", 3)

# 左：深度阶梯
steps = [
    ("2021", "Native 守护进程", "keystore2 · profcollectd", "独立进程，经 AIDL/FFI 与 C++ 交互"),
    ("2022", "协议栈 / 虚拟化", "UWB · DoH3 · AVF/Microdroid", "同一 APEX 模块内与 C++ 混合编译（cxx）"),
    ("2023", "TEE / 裸机固件", "pvmfw · Trusty TA · libavb 封装", "no_std，不依赖操作系统"),
    ("2025", "内核地址空间", "ashmem 驱动（6.12 GKI 量产）", "与 C 内核同一地址空间、同一套构建"),
    ("2026", "基带固件", "Pixel 10 基带 DNS 解析器", "可被远程攻击的最底层组件"),
]
y = Inches(1.95)
x0 = Inches(0.6)
for i, (yr, layer, mod, fence) in enumerate(steps):
    x = x0 + Inches(0.5) * i           # 阶梯向右推进
    yy = y + Inches(0.92) * i          # 向下深入
    c = RGBColor(0xF0, 0x8A, 0x4B) if i < 2 else (ACCENT if i < 4 else RGBColor(0x7A, 0x1F, 0x06))
    card = rect(s, x, yy, Inches(5.2), Inches(0.84), c, rounded=True)
    shape_text(card, "**" + yr + " · " + layer + "**   " + mod + "\n隔离方式：" + fence, 10, WHITE)

# 右：观点解读
tb = textbox(s, Inches(8.15), Inches(1.95), Inches(4.75), Inches(5.0))
tf = tb.text_frame
notes = [
    ("“围栏”指什么", True),
    ("早期每引入一处 Rust 都配有明确的隔离措施：独立进程、APEX 模块边界、较厚的 FFI 封装层", False),
    ("隔离为什么越来越少", True),
    ("互操作工具成熟（cxx/AIDL）→ Rust 与 C++ 在同一模块内混合编译；Rust-for-Linux 成熟 → 进入内核同一地址空间；no_std 成熟 → 进入 TEE/固件/基带", False),
    ("目前剩下的隔离只有 unsafe 块", True),
    ("进程/模块级的隔离取消后，需要人工审查的代码收敛为约 4% 的 unsafe{} 块，集中在 FFI 边界", False),
    ("为什么越低的层进入越晚", True),
    ("每往下一层都要先解决该层的工程问题：no_std、代码体积、厂商构建系统、内核 API 封装", False),
]
first = True
for txt, hd in notes:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_after = Pt(3 if hd else 8)
    p.line_spacing = 1.12
    _fill_para(p, txt, 12 if hd else 10.5, ACCENT if hd else BODY, bold_all=hd)
add_footer(s, "security.googleblog.com 2021–2026 各篇（已回源）；rust-in-aosp.md §1.3/§一", 3)

# ============ S4 代码量趋势 ============
s = content_slide("第 3 页 · 代码量", "总量五年从 10 万行到 695 万行；相对 C++ 的占比持续上升", 4)

# 左图：绝对量
cd = CategoryChartData()
cd.categories = ["2021-06\n官方", "2022-12\nA13 官方", "2025-11\n官方", "2026-06\n全树统计"]
cd.add_series("Rust（百万行）", (0.1, 1.5, 5.0, 6.95))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.95), Inches(5.9), Inches(4.1), cd)
ch = gf.chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = "Rust 绝对量：10万 → 150万 → 500万 → 695万行"
for r in ch.chart_title.text_frame.paragraphs[0].runs:
    _set_font(r, 12.5, True, INK)
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.00"M"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(10.5)
plot.data_labels.font.bold = True
ser = ch.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = ACCENT

# 右图：Rust/C++ 相对比例（全树统计）
cd2 = CategoryChartData()
cd2.categories = ["AOSP 14\n(2024)", "AOSP 16\n(2025)", "AOSP 17\n(2026)"]
cd2.add_series("Rust ÷ C/C++（全树）", (0.073, 0.075, 0.077))
cd2.add_series("C/C++ 占全树比例", (0.547, 0.519, 0.490))
gf2 = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(6.8), Inches(1.95), Inches(6.0), Inches(4.1), cd2)
ch2 = gf2.chart
ch2.has_legend = True
ch2.legend.include_in_layout = False
ch2.legend.font.size = Pt(10.5)
ch2.has_title = True
ch2.chart_title.text_frame.text = "对比 C++：Rust/C++ 比值 7.3%→7.7%；C/C++ 占比 54.7%→49%"
for r in ch2.chart_title.text_frame.paragraphs[0].runs:
    _set_font(r, 12.5, True, INK)
p2 = ch2.plots[0]
p2.has_data_labels = True
p2.data_labels.number_format = '0.0%'
p2.data_labels.number_format_is_linked = False
p2.data_labels.font.size = Pt(9.5)

tb = textbox(s, Inches(0.6), Inches(6.15), Inches(12.3), Inches(0.9))
tf = tb.text_frame
_fill_para(tf.paragraphs[0], "**新代码中的占比更高**：Android 13 新增原生代码 21% 为 Rust（2022，官方）；2025-11 官方确认 Rust 净增行数已超过 C++", 12, BODY)
p = tf.add_paragraph()
_fill_para(p, "存量 C/C++ 仍在增长（6836万→9009万行）——策略本来就是新代码用 Rust、旧代码不动；平台自己的 Rust 代码两年增长 3.6 倍（7.8万→35.8万行）", 11, MUTED)
add_footer(s, "官方博客 2021/2022/2025（已回源）；derdilla 全树统计 AOSP 14/16/17（tokei）", 4)

# ============ S5 混合编程 ============
s = content_slide("第 4 页 · 混合编程", "与 C/C++ 的互操作方案：先量化验证可行性（87% 可对接），再大规模投入", 5)

rows = [
    ("方案", "方向 / 机制", "AOSP 实战", "适用场景"),
    ("bindgen", "由 C 头文件自动生成 Rust 绑定（生成的代码是 unsafe）", "libbinder_ndk、libavb——生成代码不对外暴露，其上再封装安全 API", "已有稳定 C ABI 的库"),
    ("cxx", "C++ ↔ Rust 双向调用，类型不匹配在编译期报错", "蓝牙 gd 栈渐进迁移（bridge 代码生成写进 Android.bp）", "与 C++ 耦合较深的现有代码"),
    ("AIDL Rust 后端", "接口定义一次，自动生成 Rust/C++/Java 三种绑定", "keystore2、KeyMint HAL——接口不变、只换实现语言", "服务 / HAL 边界"),
    ("autocxx / bindgen 宏", "在 Rust 代码里直接描述要调用的 C++ 接口", "平台工具与测试组件", "调用点少的轻量集成"),
    ("cbindgen / crubit", "Rust → C 头文件 / 双向自动绑定（仍在开发）", "crubit 是 Google 100 万美元互操作专项的方向", "Rust 代码导出给 C 用"),
]
make_table(s, rows, Inches(0.6), Inches(1.95), Inches(12.2), Inches(3.3),
           [Inches(1.7), Inches(3.3), Inches(4.4), Inches(2.8)], size=11)

tb = textbox(s, Inches(0.6), Inches(5.5), Inches(12.2), Inches(1.5))
tf = tb.text_frame
_fill_para(tf.paragraphs[0], "**可行性是先用数据验证过的（2021-06 官方原文）**：最常用的 C++ 库，81% 的导出类型可以直接对接、87% 可以低成本对接；Mainline 模块为 88%/90%", 12.5, BODY)
p = tf.add_paragraph()
_fill_para(p, "**两条工程约定**：① 接口保持粗粒度——锁、句柄等状态不跨语言传递；② unsafe 集中管理——unsafe 只出现在 FFI 边界，全平台约 4% 的 Rust 代码在 unsafe 块内", 12.5, BODY)
p = tf.add_paragraph()
_fill_para(p, "结论：混合编程不是临时过渡——数千万行 C/C++ 会长期存在，互操作的质量直接决定 Rust 能进入多深的层", 12.5, INK)
add_footer(s, "security.googleblog.com 2021-06、2024-02（已回源）；G:\\aosp-scan Android.bp 实证", 5)

# ============ S6 项目实证 ============
s = content_slide("第 5 页 · 实际项目分析", "五个实际项目：Rust 的语言特性分别解决了什么具体问题", 6)

rows = [
    ("项目（所在层）", "要解决的问题", "用到的 Rust 特性及效果", "代码数据"),
    ("libbinder_rs\n③ IPC 基础库", "C 接口的引用计数、指针有效期、线程安全都只靠文档约定", "所有权 + Drop 自动配对引用计数；Send+Sync 使服务实现必须线程安全（否则编译失败）；Result 强制检查错误码", "308 处 unsafe 全部集中在 FFI 边界，业务代码近零"),
    ("KeyMint\n③④⑥ 密钥服务", "TEE 内存少且不可换页；输入不可信；泄露后果最严重", "枚举表达算法状态机并按算法限制输入大小；ZeroizeOnDrop 析构时清零密钥；try_reserve 处理分配失败；no_std 使同一份代码可编入 TEE", "23,392 行仅 27 处 unsafe（0.12%）"),
    ("蓝牙 GATT\n③ 协议栈（APEX）", "对端报文不可信 + 多连接并发，C 实现漏洞多", "报文由 PDL 生成类型化结构体，无手工指针操作；WeakBox 弱引用避免对象释放后回调；单线程 async 避免数据竞争", "10,497 行 GATT server 没有 unsafe 块"),
    ("libavb-rust\n⑥ 固件", "C 回调的数据所有权只在文档里说明；验证代码自身也解析不可信数据", "Ops<'a> 用生命周期参数让编译器检查 C 回调的所有权约定；Descriptor<'a> 零拷贝解析，Unknown 变体保证向前兼容", "95 处 unsafe 每处都带 SAFETY 注释"),
    ("libprefetch\n③ init 启动路径", "开机性能敏感；没有安全压力仍选了 Rust", "RAII 自动配对日志与句柄；Result 使错误路径不可忽略；直接用 crates.io 现成库（nix/serde/lru_cache）", "全模块仅 1 个 unsafe 函数（系统调用）"),
]
make_table(s, rows, Inches(0.6), Inches(1.95), Inches(12.2), Inches(3.9),
           [Inches(2.1), Inches(2.9), Inches(4.2), Inches(3.0)], size=10.5)

tb = textbox(s, Inches(0.6), Inches(6.05), Inches(12.2), Inches(1.0))
tf = tb.text_frame
_fill_para(tf.paragraphs[0], "**两种情况**：前四个项目主要解决安全问题（整类漏洞在编译期消除）；第五个说明即使不考虑安全，Rust 在开发效率上也被工程师主动选择", 12.5, BODY)
p = tf.add_paragraph()
_fill_para(p, "官方效率数据（2025-11）：Rust 代码评审耗时少约 25%、返工轮次少约 20%、中大型变更回滚率约低 4 倍", 12.5, INK)
add_footer(s, "代表性rust库分析/ 五份独立报告（G:\\aosp-scan 源码分析）；security.googleblog.com 2025-11", 6)

prs.save(r"F:\lee_space\code\research\rust-in-android-insights.pptx")
print("saved rust-in-android-insights.pptx, slides:", len(prs.slides._sldIdLst))
